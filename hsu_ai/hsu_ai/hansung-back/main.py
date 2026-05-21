from datetime import datetime
from io import BytesIO
from typing import List, Optional
import base64
import os
import re

import chromadb
from docx import Document
from dotenv import load_dotenv
from fastapi import FastAPI
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Pt
from pydantic import BaseModel, Field

try:
    from pypdf import PdfReader
except ImportError:
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        PdfReader = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

load_dotenv()

gemini = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))

app = FastAPI(title="Hansung AI Report API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],
    allow_methods=["*"],
    allow_headers=["*"],
)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
chroma_client = chromadb.Client()


def load_data(collection_name: str, filename: str):
    collection = chroma_client.get_or_create_collection(name=collection_name)
    filepath = os.path.join(BASE_DIR, filename)
    if not os.path.exists(filepath):
        return collection

    with open(filepath, "r", encoding="utf-8") as file:
        lines = [line.strip() for line in file.readlines() if line.strip()]

    if lines and len(collection.get()["ids"]) == 0:
        collection.add(documents=lines, ids=[f"doc_{index}" for index in range(len(lines))])

    return collection


collections = {
    "report": load_data("report", "hansung_data.txt"),
    "sojungdae": load_data("sojungdae", "소중대_data.txt"),
    "ai_university": load_data("ai-university", "aiuniv_data.txt"),
    "job_plus": load_data("job-plus", "jobplus_data.txt"),
}


class AttachmentSource(BaseModel):
    name: str
    mimeType: Optional[str] = None
    type: Optional[str] = None
    text: Optional[str] = None
    data: Optional[str] = None


class EmailSource(BaseModel):
    subject: str
    from_: Optional[str] = Field(default=None, alias="from")
    date: Optional[str] = None
    body: Optional[str] = None
    attachments: List[AttachmentSource] = Field(default_factory=list)


class UploadedFile(BaseModel):
    name: str
    type: Optional[str] = None
    text: Optional[str] = None
    data: Optional[str] = None


class WorkRequest(BaseModel):
    work_content: str = ""
    target_committee: str = "자료 기반 보고서"
    meeting_focus: str = "요약 보고서"
    email_sources: List[EmailSource] = Field(default_factory=list)
    uploaded_files: List[UploadedFile] = Field(default_factory=list)


class SaveRequest(BaseModel):
    content: str
    filename: str


def limit_text(text: Optional[str], max_length: int = 12000) -> str:
    if not text:
        return ""
    return text[:max_length]


def search_data(query: str) -> str:
    collection = collections["report"]
    if not query.strip():
        return "직접 입력된 업무 내용이 없어 기본 참고자료를 보조 자료로만 활용합니다."

    results = collection.query(query_texts=[query], n_results=4)
    if results["documents"] and results["documents"][0]:
        return "\n".join(results["documents"][0])
    return "관련 참고자료 없음"


def decode_file_data(data: str) -> bytes:
    normalized = data.strip()
    if "," in normalized:
        normalized = normalized.split(",", 1)[1]
    padding = "=" * (-len(normalized) % 4)
    return base64.urlsafe_b64decode(normalized + padding)


def extract_pdf_text(file_bytes: bytes) -> str:
    if PdfReader is None:
        return "(PDF 본문 추출 실패: pypdf 또는 PyPDF2가 설치되어 있지 않습니다.)"

    reader = PdfReader(BytesIO(file_bytes))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages).strip() or "(PDF에서 추출된 텍스트가 없습니다. 스캔본이면 OCR이 필요합니다.)"


def is_extraction_failure(text: str) -> bool:
    return (
        not text.strip()
        or "추출된 텍스트가 없습니다" in text
        or "본문 추출 실패" in text
        or "OCR이 필요합니다" in text
    )


def extract_docx_text(file_bytes: bytes) -> str:
    doc = Document(BytesIO(file_bytes))
    paragraphs = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
    table_rows = []
    for table in doc.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                table_rows.append(" | ".join(values))
    return "\n".join(paragraphs + table_rows).strip() or "(Word 문서에서 추출된 텍스트가 없습니다.)"


def extract_xlsx_text(file_bytes: bytes) -> str:
    if openpyxl is None:
        return "(Excel 본문 추출 실패: openpyxl이 설치되어 있지 않습니다.)"

    workbook = openpyxl.load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
    blocks = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) for value in row if value is not None and str(value).strip()]
            if values:
                rows.append(" | ".join(values))
        if rows:
            blocks.append(f"[{sheet.title}]\n" + "\n".join(rows))
    workbook.close()
    return "\n\n".join(blocks).strip() or "(Excel 문서에서 추출된 텍스트가 없습니다.)"


def extract_pptx_text(file_bytes: bytes) -> str:
    prs = Presentation(BytesIO(file_bytes))
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"[슬라이드 {index}]\n" + "\n".join(texts))
    return "\n\n".join(slides).strip() or "(PPT 문서에서 추출된 텍스트가 없습니다.)"


def extract_document_text(name: str, mime_type: Optional[str], data: Optional[str]) -> str:
    if not data:
        return ""

    try:
        file_bytes = decode_file_data(data)
        lower_name = name.lower()
        mime = mime_type or ""

        if lower_name.endswith(".pdf") or mime == "application/pdf":
            return extract_pdf_text(file_bytes)
        if lower_name.endswith(".docx") or "wordprocessingml" in mime:
            return extract_docx_text(file_bytes)
        if lower_name.endswith(".xlsx") or "spreadsheetml" in mime:
            return extract_xlsx_text(file_bytes)
        if lower_name.endswith(".pptx") or "presentationml" in mime:
            return extract_pptx_text(file_bytes)

        return "(지원하지 않는 문서 형식입니다. 파일명과 형식만 참고합니다.)"
    except Exception as exc:
        return f"(문서 본문 추출 실패: {exc})"


def source_text(name: str, mime_type: Optional[str], text: Optional[str], data: Optional[str]) -> str:
    if text:
        return text
    return extract_document_text(name, mime_type, data)


def is_pdf_source(name: str, mime_type: Optional[str]) -> bool:
    mime = mime_type or ""
    return name.lower().endswith(".pdf") or mime == "application/pdf"


def collect_gemini_file_parts(email_sources: List[EmailSource], uploaded_files: List[UploadedFile]):
    parts = []

    def append_pdf_part(name: str, mime_type: Optional[str], data: Optional[str], extracted_text: str):
        if not data or not is_pdf_source(name, mime_type):
            return
        if not is_extraction_failure(extracted_text):
            return
        try:
            file_bytes = decode_file_data(data)
            parts.append(f"\n[PDF 원본 첨부: {name}]\n위 PDF는 텍스트 추출이 어려워 원본 파일을 함께 제공합니다. 가능한 범위에서 내용을 읽고 요약에 반영하세요.")
            parts.append(types.Part.from_bytes(data=file_bytes, mime_type="application/pdf"))
        except Exception as exc:
            parts.append(f"\n[PDF 원본 첨부 실패: {name}] {exc}")

    for email in email_sources:
        for file in email.attachments:
            mime_type = file.mimeType or file.type
            extracted = source_text(file.name, mime_type, file.text, file.data)
            append_pdf_part(file.name, mime_type, file.data, extracted)

    for file in uploaded_files:
        extracted = source_text(file.name, file.type, file.text, file.data)
        append_pdf_part(file.name, file.type, file.data, extracted)

    return parts


def format_email_sources(email_sources: List[EmailSource]) -> str:
    if not email_sources:
        return "선택된 Gmail 참고자료 없음"

    blocks = []
    for index, email in enumerate(email_sources, start=1):
        attachments = "\n".join(
            [
                f"  - {file.name} ({file.mimeType or file.type or 'unknown'}): "
                f"{limit_text(source_text(file.name, file.mimeType or file.type, file.text, file.data), 5000)}"
                for file in email.attachments
            ]
        )
        blocks.append(
            f"""[메일 {index}]
제목: {email.subject}
발신자: {email.from_ or ''}
일자: {email.date or ''}
본문:
{limit_text(email.body)}
첨부:
{attachments or '첨부 없음'}"""
        )
    return "\n\n".join(blocks)


def format_uploaded_files(uploaded_files: List[UploadedFile]) -> str:
    if not uploaded_files:
        return "사용자 첨부자료 없음"

    return "\n\n".join(
        [
            f"""[첨부자료 {index}]
파일명: {file.name}
형식: {file.type or 'unknown'}
내용:
{limit_text(source_text(file.name, file.type, file.text, file.data))}"""
            for index, file in enumerate(uploaded_files, start=1)
        ]
    )


def strip_code_fence(content: str) -> str:
    cleaned = content.strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:markdown|md)?\s*", "", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\s*```$", "", cleaned)
    return cleaned.strip()


def clean_markdown_text(text: str) -> str:
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"__(.*?)__", r"\1", text)
    text = re.sub(r"`([^`]*)`", r"\1", text)
    return text.strip()


def markdown_lines(content: str) -> List[str]:
    return [line.rstrip() for line in strip_code_fence(content).splitlines()]


def add_word_paragraph(doc: Document, line: str):
    stripped = line.strip()
    if not stripped:
        doc.add_paragraph("")
        return

    heading = re.match(r"^(#{1,4})\s+(.+)$", stripped)
    if heading:
        level = min(len(heading.group(1)), 3)
        doc.add_heading(clean_markdown_text(heading.group(2)), level=level)
        return

    bullet = re.match(r"^[-*]\s+(.+)$", stripped)
    if bullet:
        doc.add_paragraph(clean_markdown_text(bullet.group(1)), style="List Bullet")
        return

    numbered = re.match(r"^\d+[.)]\s+(.+)$", stripped)
    if numbered:
        doc.add_paragraph(clean_markdown_text(numbered.group(1)), style="List Number")
        return

    doc.add_paragraph(clean_markdown_text(stripped))


def markdown_to_ppt_sections(content: str):
    lines = markdown_lines(content)
    title = "AI 보고서"
    sections = []
    current = None

    for raw_line in lines:
        line = raw_line.strip()
        if not line:
            continue

        h1 = re.match(r"^#\s+(.+)$", line)
        h2 = re.match(r"^#{2,3}\s+(.+)$", line)
        bullet = re.match(r"^[-*]\s+(.+)$", line)
        numbered = re.match(r"^\d+[.)]\s+(.+)$", line)

        if h1 and title == "AI 보고서":
            title = clean_markdown_text(h1.group(1))
            continue
        if h2:
            if current:
                sections.append(current)
            current = {"title": clean_markdown_text(h2.group(1)), "items": []}
            continue
        if bullet or numbered:
            item = bullet.group(1) if bullet else numbered.group(1)
            if current is None:
                current = {"title": "주요 내용", "items": []}
            current["items"].append(clean_markdown_text(item))
            continue

        if current is None:
            current = {"title": "개요", "items": []}
        current["items"].append(clean_markdown_text(line))

    if current:
        sections.append(current)

    if not sections:
        sections = [{"title": "주요 내용", "items": [clean_markdown_text(line) for line in lines if line.strip()]}]

    return title, sections


def chunk_items(items: List[str], size: int = 6):
    for index in range(0, len(items), size):
        yield items[index : index + size]


@app.post("/generate")
def generate_report(req: WorkRequest):
    today = datetime.now().strftime("%Y년 %m월 %d일")
    related = search_data(req.work_content)
    email_section = format_email_sources(req.email_sources)
    file_section = format_uploaded_files(req.uploaded_files)
    file_parts = collect_gemini_file_parts(req.email_sources, req.uploaded_files)

    prompt = f"""당신은 메일, 첨부파일, 사용자 요청을 바탕으로 보고서 초안을 작성하는 AI 문서 작성 보좌관입니다.

목표:
- 사용자가 입력한 요청을 최우선으로 해석합니다.
- 선택된 Gmail 본문/첨부파일, 사용자가 첨부한 파일, 기본 참고자료를 요약하고 정리합니다.
- 결과물은 Word와 PPT로 변환하기 쉬운 "보고서 형식의 Markdown 초안"으로 작성합니다.
- 사용자가 PPT 생성을 요청해도, 여기서는 PPT 변환에 적합한 Markdown 보고서 구조로 작성합니다.

작성 원칙:
- 자료에 근거해 작성하고, 확인되지 않는 사실/수치/일정은 임의로 만들지 않습니다.
- 자료만으로 확정할 수 없는 내용은 "(확인 필요)"로 표시합니다.
- 사용자가 특정 항목을 나열하면 각 항목을 보고서 섹션 또는 하위 항목으로 정리합니다.
- 문체는 공식 보고서에 맞게 간결하고 명확하게 작성합니다.
- 캘린더 일정은 사용하지 않습니다.
- 출력은 반드시 Markdown만 사용합니다.
- 첫 줄은 "# 보고서 제목" 형식으로 작성합니다.
- 주요 섹션은 "## 1. 개요", "## 2. 주요 내용"처럼 2단계 제목을 사용합니다.
- 각 섹션의 세부 내용은 bullet list로 작성합니다.
- 불필요한 설명, 코드블록, 서문, 결론 멘트는 출력하지 않습니다.

[작성일]
{today}

[보고서 성격]
{req.meeting_focus}

[사용자 요청]
{req.work_content or '직접 입력 없음'}

[기본 참고자료]
{related}

[선택된 Gmail 참고자료]
{email_section}

[사용자 첨부자료]
{file_section}

권장 구조:
# 적절한 보고서 제목

## 1. 개요
- 보고 목적과 대상 자료 요약

## 2. 주요 내용 요약
- 메일/첨부파일/사용자 요청에서 확인되는 핵심 내용

## 3. 세부 정리 내용
- 사용자가 요청한 항목별 정리

## 4. 검토 및 확인 필요 사항
- 자료 부족, PDF 텍스트 미추출, 추가 확인 필요 항목

## 5. 후속 조치
- 다음 단계, 정리 방향, 제출 또는 보고 준비 사항
"""

    contents = [prompt, *file_parts] if file_parts else prompt
    response = gemini.models.generate_content(model="gemini-2.5-flash-lite", contents=contents)
    return {"draft": strip_code_fence(response.text)}


@app.post("/save-word")
def save_word(req: SaveRequest):
    doc = Document()
    for line in markdown_lines(req.content):
        add_word_paragraph(doc, line)

    filepath = os.path.join(BASE_DIR, f"output_{datetime.now().strftime('%Y%m%d%H%M%S')}.docx")
    doc.save(filepath)
    return FileResponse(
        path=filepath,
        filename=req.filename,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )


@app.post("/save-ppt")
def save_ppt(req: SaveRequest):
    title, sections = markdown_to_ppt_sections(req.content)
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = datetime.now().strftime("%Y년 %m월 %d일")

    for section in sections:
        item_chunks = list(chunk_items(section["items"], 6)) or [[]]
        for chunk_index, items in enumerate(item_chunks):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            suffix = f" ({chunk_index + 1})" if len(item_chunks) > 1 else ""
            slide.shapes.title.text = f"{section['title']}{suffix}"[:60]
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            if not items:
                text_frame.text = "내용 없음"
                continue
            for item_index, item in enumerate(items):
                paragraph = text_frame.paragraphs[0] if item_index == 0 else text_frame.add_paragraph()
                paragraph.text = item[:240]
                paragraph.level = 0
                paragraph.font.size = Pt(18)

    filepath = os.path.join(BASE_DIR, f"output_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx")
    prs.save(filepath)
    return FileResponse(
        path=filepath,
        filename=req.filename,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
