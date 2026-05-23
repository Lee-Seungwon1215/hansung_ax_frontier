from datetime import datetime
from io import BytesIO
from typing import List, Optional
import base64
import json
import os
import re
import sqlite3
from string import Template
from urllib.parse import quote

import chromadb
from docx import Document
from dotenv import load_dotenv
from fastapi import FastAPI
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field

try:
    from pypdf import PdfReader
except ImportError:
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        PdfReader = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

load_dotenv()

gemini = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))

app = FastAPI(title="Hansung AI Report API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "http://127.0.0.1:3000",
    ],
    allow_methods=["*"],
    allow_headers=["*"],
)

APP_DIR = os.path.dirname(os.path.abspath(__file__))
BASE_DIR = os.path.dirname(APP_DIR)
DB_PATH = os.path.join(BASE_DIR, "hansung_demo.db")
TEMPLATE_DIR = os.path.join(BASE_DIR, "templates")
RULES_DIR = os.path.join(BASE_DIR, "rules")
chroma_client = chromadb.Client()


def load_data(collection_name: str, filename: str):
    collection = chroma_client.get_or_create_collection(name=collection_name)
    filepath = os.path.join(BASE_DIR, filename)
    if not os.path.exists(filepath):
        return collection

    with open(filepath, "r", encoding="utf-8") as file:
        lines = [line.strip() for line in file.readlines() if line.strip()]

    if lines and len(collection.get()["ids"]) == 0:
        collection.add(documents=lines, ids=[f"doc_{index}" for index in range(len(lines))])

    return collection


RULE_DOCUMENTS = [
    ("3-1-92 AI·SW교육단규정", "3-1-92_AI_SW_education_group.pdf"),
    ("3-1-4 위임전결규정", "3-1-4_delegated_authority.pdf"),
    ("3-1-84 공유·협력사업 지원에 관한 규정", "3-1-84_collaboration_support.pdf"),
    ("3-1-9 교직원보수규정", "3-1-9_staff_compensation.pdf"),
    ("3-1-3 사무분장규정", "3-1-3_duties.pdf"),
    ("2-1-1 한성대학교 학칙", "2-1-1_school_regulations.pdf"),
    ("3-1-2 직제규정", "3-1-2_organization.pdf"),
]


def split_rule_text(text: str, chunk_size: int = 900, overlap: int = 120) -> List[str]:
    normalized = re.sub(r"\n{3,}", "\n\n", text).strip()
    article_chunks = re.split(r"(?=\n?\s*제\s*\d+\s*조)", normalized)
    chunks = []
    for block in article_chunks:
        block = re.sub(r"\s+", " ", block).strip()
        if not block:
            continue
        if len(block) <= chunk_size:
            chunks.append(block)
            continue
        start = 0
        while start < len(block):
            chunks.append(block[start : start + chunk_size])
            start += chunk_size - overlap
    return chunks


def load_rule_collection():
    collection = chroma_client.get_or_create_collection(name="hansung-rules")
    if len(collection.get()["ids"]) > 0:
        return collection
    if PdfReader is None:
        return collection

    documents = []
    metadatas = []
    ids = []
    for rule_title, filename in RULE_DOCUMENTS:
        filepath = os.path.join(RULES_DIR, filename)
        if not os.path.exists(filepath):
            continue
        try:
            reader = PdfReader(filepath)
            text = "\n".join([page.extract_text() or "" for page in reader.pages])
        except Exception:
            continue
        for index, chunk in enumerate(split_rule_text(text)):
            documents.append(chunk)
            metadatas.append({"rule": rule_title, "filename": filename})
            ids.append(f"rule_{len(ids)}")

    if documents:
        collection.add(documents=documents, metadatas=metadatas, ids=ids)
    return collection


collections = {
    "report": load_data("report", "hansung_data.txt"),
    "sojungdae": load_data("sojungdae", "소중대_data.txt"),
    "ai_university": load_data("ai-university", "aiuniv_data.txt"),
    "rules": load_rule_collection(),
}


class AttachmentSource(BaseModel):
    name: str
    mimeType: Optional[str] = None
    type: Optional[str] = None
    text: Optional[str] = None
    data: Optional[str] = None


class EmailSource(BaseModel):
    subject: str
    from_: Optional[str] = Field(default=None, alias="from")
    date: Optional[str] = None
    body: Optional[str] = None
    attachments: List[AttachmentSource] = Field(default_factory=list)


class UploadedFile(BaseModel):
    name: str
    type: Optional[str] = None
    text: Optional[str] = None
    data: Optional[str] = None


class WorkRequest(BaseModel):
    work_content: str = ""
    target_committee: str = "자료 기반 보고서"
    meeting_focus: str = "요약 보고서"
    email_sources: List[EmailSource] = Field(default_factory=list)
    uploaded_files: List[UploadedFile] = Field(default_factory=list)


class SaveRequest(BaseModel):
    content: str
    filename: str


class WeeklyReportCreateRequest(BaseModel):
    title: str
    period: str = ""
    source_approval_ids: List[str] = Field(default_factory=list)
    draft_content: str
    final_content: Optional[str] = None
    status: str = "초안"
    created_by: str = "SW중심대학사업단 운영위원회 직원"


class WeeklyReportComposeRequest(BaseModel):
    approval_ids: List[str] = Field(default_factory=list)


class PromptTemplateUpdateRequest(BaseModel):
    purpose: str
    sections: List[str]
    rules: List[str]
    is_active: bool = True


class AiPromptUpdateRequest(BaseModel):
    content: str
    description: str = ""
    is_active: bool = True


class ApprovalCreateRequest(BaseModel):
    title: str
    category: str
    requester: str = "SW중심대학사업단 운영위원회"
    department: str = "SW중심대학사업단"
    amount: Optional[str] = None
    content: str
    attachments: List[UploadedFile] = Field(default_factory=list)


class ApprovalActionRequest(BaseModel):
    action: str
    comment: str = ""
    actor: str = "SW중심대학사업단 운영위원회 직원"


MOCK_CURRENT_USER = {
    "id": "sw-committee-staff-01",
    "name": "운영위원회 직원",
    "department": "SW중심대학사업단",
    "role": "committee_staff",
    "permissions": ["approval:create", "approval:review", "report:create", "report:finalize"],
}


def db_connection():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def init_db():
    with db_connection() as conn:
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS approvals (
                id TEXT PRIMARY KEY,
                payload TEXT NOT NULL,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS weekly_reports (
                id TEXT PRIMARY KEY,
                title TEXT NOT NULL,
                period TEXT NOT NULL,
                source_approval_ids TEXT NOT NULL,
                draft_content TEXT NOT NULL,
                final_content TEXT NOT NULL,
                status TEXT NOT NULL,
                created_by TEXT NOT NULL,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS prompt_templates (
                report_type TEXT PRIMARY KEY,
                purpose TEXT NOT NULL,
                sections TEXT NOT NULL,
                rules TEXT NOT NULL,
                is_active INTEGER NOT NULL DEFAULT 1,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS ai_prompts (
                prompt_key TEXT PRIMARY KEY,
                content TEXT NOT NULL,
                description TEXT NOT NULL,
                is_active INTEGER NOT NULL DEFAULT 1,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL
            )
            """
        )
        conn.execute(
            """
            CREATE TABLE IF NOT EXISTS generated_files (
                id TEXT PRIMARY KEY,
                filename TEXT NOT NULL,
                file_type TEXT NOT NULL,
                media_type TEXT NOT NULL,
                content BLOB NOT NULL,
                source_content TEXT NOT NULL,
                created_at TEXT NOT NULL
            )
            """
        )


def db_approvals():
    with db_connection() as conn:
        rows = conn.execute("SELECT payload FROM approvals ORDER BY created_at DESC").fetchall()
    return [json.loads(row["payload"]) for row in rows]


def save_generated_file(filename: str, file_type: str, media_type: str, content: bytes, source_content: str):
    created_at = datetime.now().isoformat(timespec="seconds")
    file_id = f"GF-{datetime.now().strftime('%Y%m%d%H%M%S%f')}"
    with db_connection() as conn:
        conn.execute(
            """
            INSERT INTO generated_files (
                id, filename, file_type, media_type, content, source_content, created_at
            )
            VALUES (?, ?, ?, ?, ?, ?, ?)
            """,
            (file_id, filename, file_type, media_type, content, source_content, created_at),
        )
    return {
        "id": file_id,
        "filename": filename,
        "file_type": file_type,
        "media_type": media_type,
        "created_at": created_at,
    }


def generated_file_response(content: bytes, filename: str, media_type: str):
    encoded_filename = quote(filename)
    return StreamingResponse(
        BytesIO(content),
        media_type=media_type,
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{encoded_filename}"},
    )


def save_db_approval(approval: dict):
    now = datetime.now().isoformat(timespec="seconds")
    with db_connection() as conn:
        conn.execute(
            """
            INSERT INTO approvals (id, payload, created_at, updated_at)
            VALUES (?, ?, ?, ?)
            ON CONFLICT(id) DO UPDATE SET payload = excluded.payload, updated_at = excluded.updated_at
            """,
            (approval["id"], json.dumps(approval, ensure_ascii=False), now, now),
        )


def get_all_approvals():
    saved_ids = set()
    saved = []
    for approval in db_approvals():
        saved_ids.add(approval["id"])
        saved.append(approval)
    return saved + [approval for approval in INTRANET_APPROVALS if approval["id"] not in saved_ids]


def update_approval_payload(approval_id: str, updated: dict) -> bool:
    for index, approval in enumerate(INTRANET_APPROVALS):
        if approval["id"] == approval_id:
            INTRANET_APPROVALS[index] = updated
            return True

    for approval in db_approvals():
        if approval["id"] == approval_id:
            save_db_approval(updated)
            return True

    return False


def seed_prompt_templates():
    now = datetime.now().isoformat(timespec="seconds")
    with db_connection() as conn:
        count = conn.execute("SELECT COUNT(*) FROM prompt_templates").fetchone()[0]
        if count:
            return
        for report_type, template in REPORT_TYPE_TEMPLATES.items():
            conn.execute(
                """
                INSERT INTO prompt_templates (
                    report_type, purpose, sections, rules, is_active, created_at, updated_at
                )
                VALUES (?, ?, ?, ?, 1, ?, ?)
                """,
                (
                    report_type,
                    template["purpose"],
                    json.dumps(template["sections"], ensure_ascii=False),
                    json.dumps(template["rules"], ensure_ascii=False),
                    now,
                    now,
                ),
            )


def get_prompt_template(report_type: str):
    with db_connection() as conn:
        row = conn.execute(
            "SELECT * FROM prompt_templates WHERE report_type = ? AND is_active = 1",
            (report_type,),
        ).fetchone()
        if row is None:
            row = conn.execute(
                "SELECT * FROM prompt_templates WHERE report_type = ? AND is_active = 1",
                ("요약 보고서",),
            ).fetchone()
    if row is None:
        return REPORT_TYPE_TEMPLATES["요약 보고서"]
    return {
        "purpose": row["purpose"],
        "sections": json.loads(row["sections"]),
        "rules": json.loads(row["rules"]),
    }


def list_prompt_templates():
    with db_connection() as conn:
        rows = conn.execute("SELECT * FROM prompt_templates ORDER BY report_type").fetchall()
    return [
        {
            "report_type": row["report_type"],
            "purpose": row["purpose"],
            "sections": json.loads(row["sections"]),
            "rules": json.loads(row["rules"]),
            "is_active": bool(row["is_active"]),
            "created_at": row["created_at"],
            "updated_at": row["updated_at"],
        }
        for row in rows
    ]


def seed_ai_prompts():
    now = datetime.now().isoformat(timespec="seconds")
    defaults = {
        "base_report_prompt": {
            "content": DEFAULT_BASE_REPORT_PROMPT,
            "description": "AI 보고서 생성 기본 프롬프트",
        },
        "weekly_report_request_prompt": {
            "content": DEFAULT_WEEKLY_REPORT_REQUEST_PROMPT,
            "description": "선택 결재건 기반 주간보고서 요청내용 자동생성 프롬프트",
        },
    }
    with db_connection() as conn:
        for prompt_key, prompt in defaults.items():
            existing = conn.execute(
                "SELECT prompt_key FROM ai_prompts WHERE prompt_key = ?",
                (prompt_key,),
            ).fetchone()
            if existing:
                continue
            conn.execute(
                """
                INSERT INTO ai_prompts (
                    prompt_key, content, description, is_active, created_at, updated_at
                )
                VALUES (?, ?, ?, 1, ?, ?)
                """,
                (
                    prompt_key,
                    prompt["content"],
                    prompt["description"],
                    now,
                    now,
                ),
            )


def ensure_ai_prompt_defaults():
    now = datetime.now().isoformat(timespec="seconds")
    with db_connection() as conn:
        rows = conn.execute("SELECT prompt_key, content FROM ai_prompts").fetchall()
        for row in rows:
            if row["prompt_key"] == "base_report_prompt" and "$reference_section" not in row["content"]:
                conn.execute(
                    """
                    UPDATE ai_prompts
                    SET content = ?, description = ?, updated_at = ?
                    WHERE prompt_key = ?
                    """,
                    (
                        DEFAULT_BASE_REPORT_PROMPT,
                        "AI 보고서 생성 기본 프롬프트",
                        now,
                        "base_report_prompt",
                    ),
                )


def get_ai_prompt(prompt_key: str) -> str:
    with db_connection() as conn:
        row = conn.execute(
            "SELECT content FROM ai_prompts WHERE prompt_key = ? AND is_active = 1",
            (prompt_key,),
        ).fetchone()
    if row is None:
        if prompt_key == "weekly_report_request_prompt":
            return DEFAULT_WEEKLY_REPORT_REQUEST_PROMPT
        return DEFAULT_BASE_REPORT_PROMPT
    return row["content"]


def list_ai_prompts():
    with db_connection() as conn:
        rows = conn.execute("SELECT * FROM ai_prompts ORDER BY prompt_key").fetchall()
    return [
        {
            "prompt_key": row["prompt_key"],
            "content": row["content"],
            "description": row["description"],
            "is_active": bool(row["is_active"]),
            "created_at": row["created_at"],
            "updated_at": row["updated_at"],
        }
        for row in rows
    ]


def attachment_payload(file: UploadedFile):
    mime_type = file.type or "unknown"
    extracted = file.text or source_text(file.name, mime_type, None, file.data)
    return {
        "name": file.name,
        "type": mime_type,
        "text": limit_text(extracted, 5000) if extracted else "",
        "hasContent": bool(extracted),
    }


init_db()


INTRANET_NOTICES = [
    {
        "id": "N-2026-001",
        "title": "2026학년도 SW중심대학사업 연간 실행계획 심의 일정 안내",
        "department": "SW중심대학사업단",
        "date": "2026-05-20",
        "priority": "중요",
        "summary": "센터별 실행계획, 예산 배정안, KPI 목표치를 운영위원회 심의 안건으로 제출해 주시기 바랍니다.",
    },
    {
        "id": "N-2026-002",
        "title": "산학협력 프로젝트 및 기업연계 인턴십 운영계획 취합",
        "department": "산학협력센터",
        "date": "2026-05-18",
        "priority": "안내",
        "summary": "협력 기업별 프로젝트 주제, 참여 학생 수, 예산 소요 내역을 사업단 양식에 맞춰 등록해 주세요.",
    },
    {
        "id": "N-2026-003",
        "title": "장학금 및 연구비 지급 기준 개정안 검토 요청",
        "department": "교육혁신센터",
        "date": "2026-05-15",
        "priority": "검토",
        "summary": "SW 전공·융합전공 학생 장학금과 산학 프로젝트 참여 연구원 지원금 기준 개정안을 확인해 주세요.",
    },
]


INTRANET_APPROVALS = [
    {
        "id": "A-2026-1042",
        "title": "2026년 SW중심대학사업 연간 사업계획 및 예산 배정안",
        "category": "심의/의결서",
        "requester": "사업운영팀",
        "department": "SW중심대학사업단",
        "submittedAt": "2026-05-21 09:40",
        "status": "진행중",
        "amount": "3,850,000,000원",
        "currentApprover": "운영위원장",
        "line": ["사업운영팀", "사업단장", "운영위원회"],
        "workType": "사업계획/예산",
        "outputType": "심의/의결서",
        "content": "정부 지원금과 교비 매칭 예산을 교육혁신, 산학협력, 가치확산, 성과관리 영역별로 배정하는 안건입니다.",
        "result": "연간 실행계획과 센터별 예산 배정안을 운영위원회 심의 안건으로 상정했습니다.",
        "followUp": "운영위원회 의결 결과에 따라 과학기술정보통신부 제출용 사업계획서에 반영해야 합니다.",
        "attachments": [
            {"name": "2026_SW중심대학사업_연간사업계획서.docx", "type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"},
            {"name": "센터별_예산배정안.xlsx", "type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"},
        ],
    },
    {
        "id": "A-2026-1041",
        "title": "SW중심대학사업 예산 변경 요청서",
        "category": "결재서",
        "requester": "사업운영팀",
        "department": "SW중심대학사업단",
        "submittedAt": "2026-05-20 16:30",
        "status": "진행중",
        "amount": "85,000,000원",
        "currentApprover": "사업단장",
        "line": ["사업운영팀", "사업단장", "운영위원회"],
        "workType": "사업계획/예산",
        "outputType": "결재서",
        "content": "산학 프로젝트 참여 기업 증가에 따라 기업연계 프로젝트 운영비와 학생 활동 지원비의 예산 조정을 요청합니다.",
        "result": "기업연계 프로젝트 확대에 필요한 예산 변경 사유와 조정 금액을 정리해 결재 상신했습니다.",
        "followUp": "승인 후 센터별 예산 집행 계획과 세부 산출내역을 갱신해야 합니다.",
        "attachments": [
            {"name": "예산변경_사유서.pdf", "type": "application/pdf"},
            {"name": "기업연계프로젝트_운영비_산출내역.xlsx", "type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"},
        ],
    },
    {
        "id": "A-2026-1037",
        "title": "비전공자 SW기초교육 마이크로디그리 개편 제안",
        "category": "제안서",
        "requester": "교육혁신센터",
        "department": "SW중심대학사업단",
        "submittedAt": "2026-05-19 14:12",
        "status": "결재대기",
        "amount": "-",
        "currentApprover": "교육분과위원",
        "line": ["센터장", "사업단장", "운영위원회"],
        "workType": "교육정책/산학협력",
        "outputType": "제안서",
        "content": "AI·SW 기초역량 강화를 위해 비전공자 대상 교과목과 마이크로디그리 이수 체계를 조정하는 제안입니다.",
        "result": "비전공자 SW기초교육 체계 개편안을 작성하고 교육분과 검토 단계로 넘겼습니다.",
        "followUp": "학사 운영 부서와 교과목 개설 가능 여부 및 이수 기준을 추가 확인해야 합니다.",
        "attachments": [
            {"name": "비전공자_SW기초교육_개편제안서.docx", "type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"},
        ],
    },
    {
        "id": "A-2026-1035",
        "title": "SW 장학금 지급 기준 개정 회의록",
        "category": "회의록",
        "requester": "교육혁신센터",
        "department": "SW중심대학사업단",
        "submittedAt": "2026-05-18 17:20",
        "status": "결재대기",
        "amount": "-",
        "currentApprover": "교육분과위원",
        "line": ["교육혁신센터", "사업단장", "운영위원회"],
        "workType": "운영규정/지침",
        "outputType": "회의록",
        "content": "전공자, 융합전공자, 산학 프로젝트 참여 학생의 장학금 지급 대상과 차등 지급 기준을 논의한 회의 결과입니다.",
        "result": "장학금 지급 대상, 평가 기준, 차등 지급 항목에 대한 회의 내용을 회의록으로 정리했습니다.",
        "followUp": "개정 기준안을 심의/의결서로 전환하고 개인정보 포함 자료의 접근 권한을 점검해야 합니다.",
        "attachments": [
            {"name": "장학금_지급기준_개정회의록.pdf", "type": "application/pdf"},
            {"name": "장학금_지급기준_비교표.xlsx", "type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"},
        ],
    },
    {
        "id": "A-2026-1031",
        "title": "성북구 청소년 SW가치확산 프로그램 운영 제안",
        "category": "제안서",
        "requester": "가치확산센터",
        "department": "SW중심대학사업단",
        "submittedAt": "2026-05-17 10:10",
        "status": "완료",
        "amount": "24,000,000원",
        "currentApprover": "결재완료",
        "line": ["가치확산센터", "사업단장", "운영위원회"],
        "workType": "교육정책/산학협력",
        "outputType": "제안서",
        "content": "성북구 초·중·고교 대상 AI·SW 체험 교육을 확대하기 위한 프로그램 운영 계획과 예산 사용안을 제안합니다.",
        "result": "성북구 연계 가치확산 프로그램 운영 계획을 승인 완료했습니다.",
        "followUp": "참여 학교 모집 일정과 강사 배정 계획을 확정하고 성과지표에 반영해야 합니다.",
        "attachments": [
            {"name": "성북구_SW가치확산_운영제안서.docx", "type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"},
        ],
    },
    {
        "id": "A-2026-1028",
        "title": "2026년 1분기 KPI 달성 현황 및 환류 계획 보고",
        "category": "보고서",
        "requester": "성과관리팀",
        "department": "SW중심대학사업단",
        "submittedAt": "2026-05-16 11:05",
        "status": "완료",
        "amount": "-",
        "currentApprover": "결재완료",
        "line": ["성과관리팀", "사업단장", "운영위원회"],
        "workType": "성과지표/환류",
        "outputType": "보고서",
        "content": "취업률, 인턴십 참여율, 오픈소스 활동, 가치확산 실적을 점검하고 미달 지표의 개선 계획을 보고합니다.",
        "result": "1분기 KPI 달성 현황을 점검하고 미달 지표의 개선 방향을 보고서로 확정했습니다.",
        "followUp": "인턴십 참여율과 오픈소스 활동 실적 보완 계획을 다음 운영위원회 안건에 포함해야 합니다.",
        "attachments": [
            {"name": "1분기_KPI_달성현황_보고서.pptx", "type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"},
            {"name": "성과지표_원자료.xlsx", "type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"},
        ],
    },
]


WORK_CATEGORIES = [
    {
        "name": "사업계획/예산",
        "description": "연간 실행계획, 정부지원금·교비 예산 편성, 예산 변경 및 집행 타당성 검토",
    },
    {
        "name": "운영규정/지침",
        "description": "사업단 운영 규정, 장학금·연구비 지급 기준, 자산·보안 지침 제·개정",
    },
    {
        "name": "교육정책/산학협력",
        "description": "SW 전공·융합 교육과정 개편, 마이크로디그리, 기업 프로젝트 및 인턴십 운영 방향",
    },
    {
        "name": "성과지표/환류",
        "description": "취업률, 인턴십 참여율, 오픈소스, 가치확산 실적 점검 및 개선 계획 관리",
    },
]


WEEKLY_REPORT_TEMPLATE = {
    "title": "SW중심대학사업단 운영위원회 주간업무보고",
    "sections": [
        "1. 금주 업무 총괄",
        "2. 주요 결재 및 심의 처리 현황",
        "3. 사업계획·예산 관련 업무",
        "4. 운영 규정·지침 관련 업무",
        "5. 교육정책·산학협력 관련 업무",
        "6. 성과지표·환류 관련 업무",
        "7. 차주 예정 업무 및 확인 필요 사항",
    ],
    "writingRules": [
        "결재 이력에 있는 사실만 근거로 작성",
        "금액, 일정, 승인 여부는 원문 값이 없으면 확인 필요로 표시",
        "운영위원회 보고에 맞는 공식 문체 사용",
        "각 업무별 결과와 후속 조치를 분리해 작성",
    ],
}


INTRANET_DOCUMENTS = [
    {
        "id": "D-001",
        "title": "안건서 양식",
        "type": "docx",
        "owner": "사업운영팀",
        "updatedAt": "2026-05-12",
        "templateFile": "agenda_ai_coding_camp_plan.docx",
    },
    {
        "id": "D-002",
        "title": "운영위원회 회의록 양식",
        "type": "docx",
        "owner": "사업운영팀",
        "updatedAt": "2026-04-30",
        "templateFile": "committee_minutes.docx",
    },
    {
        "id": "D-003",
        "title": "심의·의결서 - AISW 산학연계 프로젝트 운영 계획(안)",
        "type": "docx",
        "owner": "SW중심대학사업단",
        "updatedAt": "2026-04-22",
        "templateFile": "deliberation_aisw_project_plan.docx",
    },
    {
        "id": "D-004",
        "title": "운영위원회 결과보고서 양식",
        "type": "docx",
        "owner": "성과관리팀",
        "updatedAt": "2026-04-18",
        "templateFile": "committee_result_report.docx",
    },
    {
        "id": "D-005",
        "title": "결재서 - 운영계획 양식",
        "type": "docx",
        "owner": "사업운영팀",
        "updatedAt": "2026-04-18",
        "templateFile": "approval_operation_plan.docx",
    },
]


INTRANET_CONTACTS = [
    {"name": "사업단장", "department": "SW중심대학사업단", "role": "사업 총괄 및 운영위원회 상정", "email": "sw-director@hansung.ac.kr", "phone": "02-760-4100"},
    {"name": "사업운영팀", "department": "SW중심대학사업단", "role": "예산·결재·위원회 운영", "email": "sw-office@hansung.ac.kr", "phone": "02-760-4101"},
    {"name": "교육혁신센터", "department": "SW중심대학사업단", "role": "교육과정·장학금 기준", "email": "sw-edu@hansung.ac.kr", "phone": "02-760-4102"},
    {"name": "산학협력센터", "department": "SW중심대학사업단", "role": "기업 프로젝트·인턴십", "email": "sw-industry@hansung.ac.kr", "phone": "02-760-4103"},
    {"name": "성과관리팀", "department": "SW중심대학사업단", "role": "KPI·자체평가·환류", "email": "sw-kpi@hansung.ac.kr", "phone": "02-760-4104"},
]


@app.get("/intranet/dashboard")
def intranet_dashboard():
    approvals = get_all_approvals()
    waiting = [item for item in approvals if item["status"] == "결재대기"]
    active = [item for item in approvals if item["status"] == "진행중"]
    done = [item for item in approvals if item["status"] == "완료"]
    return {
        "summary": {
            "waitingApprovals": len(waiting),
            "activeApprovals": len(active),
            "completedApprovals": len(done),
            "notices": len(INTRANET_NOTICES),
        },
        "currentUser": MOCK_CURRENT_USER,
        "notices": INTRANET_NOTICES,
        "approvals": approvals,
        "documents": INTRANET_DOCUMENTS,
        "contacts": INTRANET_CONTACTS,
        "workCategories": WORK_CATEGORIES,
        "weeklyReportTemplate": WEEKLY_REPORT_TEMPLATE,
    }


@app.get("/auth/me")
def auth_me():
    return MOCK_CURRENT_USER


@app.get("/documents/{document_id}/download")
def download_document_template(document_id: str):
    document = next((item for item in INTRANET_DOCUMENTS if item["id"] == document_id), None)
    if document is None:
        return {"ok": False, "message": "양식을 찾을 수 없습니다."}

    template_file = document.get("templateFile")
    filepath = os.path.join(TEMPLATE_DIR, template_file)
    if not template_file or not os.path.exists(filepath):
        return {"ok": False, "message": "양식 파일이 없습니다."}

    return FileResponse(
        path=filepath,
        filename=f"{document['title']}.docx",
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )


@app.get("/rules/search")
def search_rule_documents(q: str):
    return {"query": q, "results": search_rules(q)}


@app.get("/prompt-templates")
def get_prompt_templates():
    return {"templates": list_prompt_templates()}


@app.put("/prompt-templates/{report_type}")
def update_prompt_template(report_type: str, req: PromptTemplateUpdateRequest):
    now = datetime.now().isoformat(timespec="seconds")
    with db_connection() as conn:
        existing = conn.execute(
            "SELECT report_type FROM prompt_templates WHERE report_type = ?",
            (report_type,),
        ).fetchone()
        if existing:
            conn.execute(
                """
                UPDATE prompt_templates
                SET purpose = ?, sections = ?, rules = ?, is_active = ?, updated_at = ?
                WHERE report_type = ?
                """,
                (
                    req.purpose,
                    json.dumps(req.sections, ensure_ascii=False),
                    json.dumps(req.rules, ensure_ascii=False),
                    1 if req.is_active else 0,
                    now,
                    report_type,
                ),
            )
        else:
            conn.execute(
                """
                INSERT INTO prompt_templates (
                    report_type, purpose, sections, rules, is_active, created_at, updated_at
                )
                VALUES (?, ?, ?, ?, ?, ?, ?)
                """,
                (
                    report_type,
                    req.purpose,
                    json.dumps(req.sections, ensure_ascii=False),
                    json.dumps(req.rules, ensure_ascii=False),
                    1 if req.is_active else 0,
                    now,
                    now,
                ),
            )
    return {"ok": True, "report_type": report_type, "updated_at": now}


@app.get("/ai-prompts")
def get_ai_prompts():
    return {"prompts": list_ai_prompts()}


@app.put("/ai-prompts/{prompt_key}")
def update_ai_prompt(prompt_key: str, req: AiPromptUpdateRequest):
    now = datetime.now().isoformat(timespec="seconds")
    description = req.description or f"{prompt_key} 프롬프트"
    with db_connection() as conn:
        existing = conn.execute(
            "SELECT prompt_key FROM ai_prompts WHERE prompt_key = ?",
            (prompt_key,),
        ).fetchone()
        if existing:
            conn.execute(
                """
                UPDATE ai_prompts
                SET content = ?, description = ?, is_active = ?, updated_at = ?
                WHERE prompt_key = ?
                """,
                (req.content, description, 1 if req.is_active else 0, now, prompt_key),
            )
        else:
            conn.execute(
                """
                INSERT INTO ai_prompts (
                    prompt_key, content, description, is_active, created_at, updated_at
                )
                VALUES (?, ?, ?, ?, ?, ?)
                """,
                (prompt_key, req.content, description, 1 if req.is_active else 0, now, now),
            )
    return {"ok": True, "prompt_key": prompt_key, "updated_at": now}


@app.post("/intranet/approvals")
def create_intranet_approval(req: ApprovalCreateRequest):
    approvals = get_all_approvals()
    new_item = {
        "id": f"A-2026-{1043 + len(approvals)}",
        "title": req.title,
        "category": req.category,
        "requester": req.requester,
        "department": req.department,
        "submittedAt": datetime.now().strftime("%Y-%m-%d %H:%M"),
        "status": "결재대기",
        "amount": req.amount or "-",
        "currentApprover": "운영위원회 검토",
        "line": ["담당센터", "사업단장", "운영위원회"],
        "workType": "확인 필요",
        "outputType": req.category,
        "content": req.content,
        "result": "신규 결재 건으로 등록되었습니다.",
        "followUp": "담당자가 업무 분류와 후속 조치를 보완해야 합니다.",
        "attachments": [attachment_payload(file) for file in req.attachments],
        "history": [
            {
                "actor": req.requester,
                "action": "상신",
                "comment": "결재 기안이 등록되었습니다.",
                "createdAt": datetime.now().strftime("%Y-%m-%d %H:%M"),
            }
        ],
    }
    save_db_approval(new_item)
    return new_item


@app.post("/intranet/approvals/{approval_id}/action")
def action_intranet_approval(approval_id: str, req: ApprovalActionRequest):
    approvals = get_all_approvals()
    approval = next((item for item in approvals if item["id"] == approval_id), None)
    if approval is None:
        return {"ok": False, "message": "결재건을 찾을 수 없습니다."}

    action_map = {
        "submit": "결재대기",
        "review": "진행중",
        "approve": "완료",
        "reject": "반려",
    }
    next_status = action_map.get(req.action)
    if not next_status:
        return {"ok": False, "message": "지원하지 않는 결재 처리입니다."}

    approval["status"] = next_status
    approval["currentApprover"] = "결재완료" if next_status == "완료" else req.actor
    approval.setdefault("history", []).append(
        {
            "actor": req.actor,
            "action": req.action,
            "comment": req.comment,
            "createdAt": datetime.now().strftime("%Y-%m-%d %H:%M"),
        }
    )
    update_approval_payload(approval_id, approval)
    return {"ok": True, "approval": approval}


@app.post("/weekly-reports")
def create_weekly_report(req: WeeklyReportCreateRequest):
    now = datetime.now().isoformat(timespec="seconds")
    report_id = f"WR-{datetime.now().strftime('%Y%m%d%H%M%S')}"
    final_content = req.final_content if req.final_content is not None else req.draft_content
    with db_connection() as conn:
        conn.execute(
            """
            INSERT INTO weekly_reports (
                id, title, period, source_approval_ids, draft_content, final_content,
                status, created_by, created_at, updated_at
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                report_id,
                req.title,
                req.period,
                json.dumps(req.source_approval_ids, ensure_ascii=False),
                req.draft_content,
                final_content,
                req.status,
                req.created_by,
                now,
                now,
            ),
        )
    return {
        "id": report_id,
        "title": req.title,
        "period": req.period,
        "source_approval_ids": req.source_approval_ids,
        "draft_content": req.draft_content,
        "final_content": final_content,
        "status": req.status,
        "created_by": req.created_by,
        "created_at": now,
        "updated_at": now,
    }


@app.post("/weekly-reports/compose-request")
def compose_weekly_report_request(req: WeeklyReportComposeRequest):
    approvals = get_all_approvals()
    selected = [item for item in approvals if item["id"] in req.approval_ids]
    return {
        "work_content": render_weekly_report_request(
            {
                "approval_count": str(len(selected)),
                "category_text": format_work_categories(),
                "template_text": format_weekly_template(),
                "weekly_work_text": format_weekly_work_items(selected),
            }
        )
    }


@app.get("/weekly-reports")
def list_weekly_reports():
    with db_connection() as conn:
        rows = conn.execute("SELECT * FROM weekly_reports ORDER BY created_at DESC").fetchall()
    reports = []
    for row in rows:
        item = dict(row)
        item["source_approval_ids"] = json.loads(item["source_approval_ids"])
        reports.append(item)
    return {"reports": reports}


def list_generated_files():
    with db_connection() as conn:
        rows = conn.execute(
            """
            SELECT id, filename, file_type, media_type, length(content) AS size, created_at
            FROM generated_files
            ORDER BY created_at DESC
            """
        ).fetchall()
    return {"files": [dict(row) for row in rows]}


def download_generated_file(file_id: str):
    with db_connection() as conn:
        row = conn.execute("SELECT * FROM generated_files WHERE id = ?", (file_id,)).fetchone()
    if row is None:
        return {"ok": False, "message": "저장된 파일을 찾을 수 없습니다."}
    return generated_file_response(row["content"], row["filename"], row["media_type"])


def limit_text(text: Optional[str], max_length: int = 12000) -> str:
    if not text:
        return ""
    return text[:max_length]


REFERENCE_COLLECTIONS = [
    ("한성대학교 기본 정보", "report", 4),
    ("SW중심대학 사업단 정보", "sojungdae", 4),
    ("AI·SW 교육 정보", "ai_university", 3),
]


def query_collection(collection_name: str, query: str, n_results: int) -> List[str]:
    collection = collections[collection_name]
    results = collection.query(query_texts=[query], n_results=n_results)
    return results.get("documents", [[]])[0] if results.get("documents") else []


def search_reference_sources(query: str) -> str:
    if not query.strip():
        return "직접 입력된 업무 내용이 없어 기본 참고자료를 보조 자료로만 활용합니다."

    blocks = []
    for label, collection_name, n_results in REFERENCE_COLLECTIONS:
        documents = query_collection(collection_name, query, n_results)
        if documents:
            blocks.append(f"[{label}]\n" + "\n".join([doc[:900] for doc in documents]))
    return "\n\n".join(blocks) if blocks else "관련 참고자료 없음"


def search_data(query: str) -> str:
    return search_reference_sources(query)


def search_rules(query: str) -> str:
    collection = collections["rules"]
    if not query.strip():
        return "관련 규정 검색어 없음"

    results = collection.query(query_texts=[query], n_results=5)
    documents = results.get("documents", [[]])[0]
    metadatas = results.get("metadatas", [[]])[0]
    if not documents:
        return "확인된 관련 규정 없음"

    blocks = []
    for doc, metadata in zip(documents, metadatas):
        rule_name = metadata.get("rule", "규정")
        blocks.append(f"[{rule_name}]\n{doc[:1200]}")
    return "\n\n".join(blocks)


def decode_file_data(data: str) -> bytes:
    normalized = data.strip()
    if "," in normalized:
        normalized = normalized.split(",", 1)[1]
    padding = "=" * (-len(normalized) % 4)
    return base64.urlsafe_b64decode(normalized + padding)


def extract_pdf_text(file_bytes: bytes) -> str:
    if PdfReader is None:
        return "(PDF 본문 추출 실패: pypdf 또는 PyPDF2가 설치되어 있지 않습니다.)"

    reader = PdfReader(BytesIO(file_bytes))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages).strip() or "(PDF에서 추출된 텍스트가 없습니다. 스캔본이면 OCR이 필요합니다.)"


def is_extraction_failure(text: str) -> bool:
    return (
        not text.strip()
        or "추출된 텍스트가 없습니다" in text
        or "본문 추출 실패" in text
        or "OCR이 필요합니다" in text
    )


def extract_docx_text(file_bytes: bytes) -> str:
    doc = Document(BytesIO(file_bytes))
    paragraphs = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
    table_rows = []
    for table in doc.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                table_rows.append(" | ".join(values))
    return "\n".join(paragraphs + table_rows).strip() or "(Word 문서에서 추출된 텍스트가 없습니다.)"


def extract_xlsx_text(file_bytes: bytes) -> str:
    if openpyxl is None:
        return "(Excel 본문 추출 실패: openpyxl이 설치되어 있지 않습니다.)"

    workbook = openpyxl.load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
    blocks = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) for value in row if value is not None and str(value).strip()]
            if values:
                rows.append(" | ".join(values))
        if rows:
            blocks.append(f"[{sheet.title}]\n" + "\n".join(rows))
    workbook.close()
    return "\n\n".join(blocks).strip() or "(Excel 문서에서 추출된 텍스트가 없습니다.)"


def extract_pptx_text(file_bytes: bytes) -> str:
    prs = Presentation(BytesIO(file_bytes))
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"[슬라이드 {index}]\n" + "\n".join(texts))
    return "\n\n".join(slides).strip() or "(PPT 문서에서 추출된 텍스트가 없습니다.)"


def extract_document_text(name: str, mime_type: Optional[str], data: Optional[str]) -> str:
    if not data:
        return ""

    try:
        file_bytes = decode_file_data(data)
        lower_name = name.lower()
        mime = mime_type or ""

        if lower_name.endswith(".pdf") or mime == "application/pdf":
            return extract_pdf_text(file_bytes)
        if lower_name.endswith(".docx") or "wordprocessingml" in mime:
            return extract_docx_text(file_bytes)
        if lower_name.endswith(".xlsx") or "spreadsheetml" in mime:
            return extract_xlsx_text(file_bytes)
        if lower_name.endswith(".pptx") or "presentationml" in mime:
            return extract_pptx_text(file_bytes)

        return "(지원하지 않는 문서 형식입니다. 파일명과 형식만 참고합니다.)"
    except Exception as exc:
        return f"(문서 본문 추출 실패: {exc})"


def source_text(name: str, mime_type: Optional[str], text: Optional[str], data: Optional[str]) -> str:
    if text:
        return text
    return extract_document_text(name, mime_type, data)


def is_pdf_source(name: str, mime_type: Optional[str]) -> bool:
    mime = mime_type or ""
    return name.lower().endswith(".pdf") or mime == "application/pdf"


def collect_gemini_file_parts(email_sources: List[EmailSource], uploaded_files: List[UploadedFile]):
    parts = []

    def append_pdf_part(name: str, mime_type: Optional[str], data: Optional[str], extracted_text: str):
        if not data or not is_pdf_source(name, mime_type):
            return
        if not is_extraction_failure(extracted_text):
            return
        try:
            file_bytes = decode_file_data(data)
            parts.append(f"\n[PDF 원본 첨부: {name}]\n위 PDF는 텍스트 추출이 어려워 원본 파일을 함께 제공합니다. 가능한 범위에서 내용을 읽고 요약에 반영하세요.")
            parts.append(types.Part.from_bytes(data=file_bytes, mime_type="application/pdf"))
        except Exception as exc:
            parts.append(f"\n[PDF 원본 첨부 실패: {name}] {exc}")

    for email in email_sources:
        for file in email.attachments:
            mime_type = file.mimeType or file.type
            extracted = source_text(file.name, mime_type, file.text, file.data)
            append_pdf_part(file.name, mime_type, file.data, extracted)

    for file in uploaded_files:
        extracted = source_text(file.name, file.type, file.text, file.data)
        append_pdf_part(file.name, file.type, file.data, extracted)

    return parts


def format_email_sources(email_sources: List[EmailSource]) -> str:
    if not email_sources:
        return "선택된 Gmail 참고자료 없음"

    blocks = []
    for index, email in enumerate(email_sources, start=1):
        attachments = "\n".join(
            [
                f"  - {file.name} ({file.mimeType or file.type or 'unknown'}): "
                f"{limit_text(source_text(file.name, file.mimeType or file.type, file.text, file.data), 5000)}"
                for file in email.attachments
            ]
        )
        blocks.append(
            f"""[메일 {index}]
제목: {email.subject}
발신자: {email.from_ or ''}
일자: {email.date or ''}
본문:
{limit_text(email.body)}
첨부:
{attachments or '첨부 없음'}"""
        )
    return "\n\n".join(blocks)


def format_uploaded_files(uploaded_files: List[UploadedFile]) -> str:
    if not uploaded_files:
        return "사용자 첨부자료 없음"

    return "\n\n".join(
        [
            f"""[첨부자료 {index}]
파일명: {file.name}
형식: {file.type or 'unknown'}
내용:
{limit_text(source_text(file.name, file.type, file.text, file.data))}"""
            for index, file in enumerate(uploaded_files, start=1)
        ]
    )


def strip_code_fence(content: str) -> str:
    cleaned = content.strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:markdown|md)?\s*", "", cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r"\s*```$", "", cleaned)
    return cleaned.strip()


def clean_markdown_text(text: str) -> str:
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"__(.*?)__", r"\1", text)
    text = re.sub(r"`([^`]*)`", r"\1", text)
    return text.strip()


def markdown_lines(content: str) -> List[str]:
    return [line.rstrip() for line in strip_code_fence(content).splitlines()]


def add_word_paragraph(doc: Document, line: str):
    stripped = line.strip()
    if not stripped:
        doc.add_paragraph("")
        return

    heading = re.match(r"^(#{1,4})\s+(.+)$", stripped)
    if heading:
        level = min(len(heading.group(1)), 3)
        doc.add_heading(clean_markdown_text(heading.group(2)), level=level)
        return

    bullet = re.match(r"^[-*]\s+(.+)$", stripped)
    if bullet:
        try:
            doc.add_paragraph(clean_markdown_text(bullet.group(1)), style="List Bullet")
        except KeyError:
            doc.add_paragraph(f"• {clean_markdown_text(bullet.group(1))}")
        return

    numbered = re.match(r"^\d+[.)]\s+(.+)$", stripped)
    if numbered:
        try:
            doc.add_paragraph(clean_markdown_text(numbered.group(1)), style="List Number")
        except KeyError:
            doc.add_paragraph(clean_markdown_text(stripped))
        return

    doc.add_paragraph(clean_markdown_text(stripped))


def set_cell_text(cell, text: str):
    cell.text = ""
    paragraph = cell.paragraphs[0]
    run = paragraph.add_run(text)
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(9)


def extract_weekly_table_rows(content: str):
    _, sections = markdown_to_ppt_sections(content)
    rows = []
    next_week = []
    for section in sections:
        title = section["title"]
        items = section["items"]
        if "차주" in title or "향후" in title:
            next_week.extend(items)
            continue
        if not items:
            continue
        rows.append(
            {
                "category": re.sub(r"^\d+\.\s*", "", title)[:18],
                "details": "\n".join([f"• {item}" for item in items[:3]])[:700],
                "previous": "",
                "current": "\n".join([f"• {item}" for item in items[:2]])[:420],
                "note": "",
            }
        )
    return rows[:3], next_week[:5]


def fill_weekly_report_template(doc: Document, content: str):
    rows, next_week = extract_weekly_table_rows(content)

    if doc.tables:
        header_table = doc.tables[0]
        if len(header_table.rows) >= 4:
            set_cell_text(header_table.cell(0, 1), "SW중심대학사업단")
            set_cell_text(header_table.cell(1, 1), "운영위원회 담당")
            set_cell_text(header_table.cell(2, 1), "운영위원회 직원")
            set_cell_text(header_table.cell(3, 1), datetime.now().strftime("%Y년 %m월 %d일"))

    if len(doc.tables) >= 2:
        work_table = doc.tables[1]
        while len(work_table.rows) < max(4, len(rows) + 1):
            work_table.add_row()
        for index, row_data in enumerate(rows, start=1):
            row = work_table.rows[index]
            set_cell_text(row.cells[0], row_data["category"])
            set_cell_text(row.cells[1], row_data["details"])
            set_cell_text(row.cells[2], "100")
            set_cell_text(row.cells[3], row_data["previous"])
            set_cell_text(row.cells[4], row_data["current"])
            set_cell_text(row.cells[5], row_data["note"])

    if len(doc.tables) >= 3 and next_week:
        set_cell_text(doc.tables[2].cell(1, 0), "\n".join([f"• {item}" for item in next_week]))

    doc.add_page_break()
    doc.add_heading("AI 생성 보고서 상세 초안", level=1)
    for line in markdown_lines(content):
        add_word_paragraph(doc, line)


def markdown_to_ppt_sections(content: str):
    lines = markdown_lines(content)
    title = "AI 보고서"
    sections = []
    current = None

    for raw_line in lines:
        line = raw_line.strip()
        if not line:
            continue

        h1 = re.match(r"^#\s+(.+)$", line)
        h2 = re.match(r"^#{2,3}\s+(.+)$", line)
        bullet = re.match(r"^[-*]\s+(.+)$", line)
        numbered = re.match(r"^\d+[.)]\s+(.+)$", line)

        if h1 and title == "AI 보고서":
            title = clean_markdown_text(h1.group(1))
            continue
        if h2:
            if current:
                sections.append(current)
            current = {"title": clean_markdown_text(h2.group(1)), "items": []}
            continue
        if bullet or numbered:
            item = bullet.group(1) if bullet else numbered.group(1)
            if current is None:
                current = {"title": "주요 내용", "items": []}
            current["items"].append(clean_markdown_text(item))
            continue

        if current is None:
            current = {"title": "개요", "items": []}
        current["items"].append(clean_markdown_text(line))

    if current:
        sections.append(current)

    if not sections:
        sections = [{"title": "주요 내용", "items": [clean_markdown_text(line) for line in lines if line.strip()]}]

    return title, sections


def chunk_items(items: List[str], size: int = 6):
    for index in range(0, len(items), size):
        yield items[index : index + size]


def chunk_ppt_items(items: List[str], max_chars: int = 620, max_items: int = 5):
    chunk = []
    char_count = 0
    for item in items:
        item_length = len(item)
        if chunk and (len(chunk) >= max_items or char_count + item_length > max_chars):
            yield chunk
            chunk = []
            char_count = 0
        chunk.append(item)
        char_count += item_length
    if chunk:
        yield chunk


DEFAULT_BASE_REPORT_PROMPT = """$type_instruction

당신은 메일, 첨부파일, 사용자 요청을 바탕으로 보고서 초안을 작성하는 AI 문서 작성 보좌관입니다.

목표:
- 사용자가 입력한 요청을 최우선으로 해석합니다.
- 선택된 Gmail 본문/첨부파일, 사용자가 첨부한 파일, 한성대학교 참고자료와 규정자료를 요약하고 정리합니다.
- 결과물은 Word와 PPT로 변환하기 쉬운 보고서 형식의 Markdown 초안으로 작성합니다.

작성 원칙:
- 자료에 근거해 작성하고, 확인되지 않는 사실/수치/일정은 임의로 만들지 않습니다.
- 자료만으로 확정할 수 없는 내용은 "(확인 필요)"로 표시합니다.
- 규정 관련 내용은 [확인된 관련 규정 참고자료]에 있는 조항과 문장만 근거로 작성합니다.
- 관련 규정 참고자료에 없는 규정 해석, 조항 번호, 승인 권한은 임의로 작성하지 말고 "(확인 필요)"로 표시합니다.
- 한성대학교 성과, 현황, 사업단 운영 내용은 [한성대학교·사업단 참고자료]에 확인된 내용만 근거로 작성합니다.
- 주요 섹션은 [선택한 보고서 유형별 작성 지시]의 섹션 제목과 순서를 반드시 사용합니다.
- 각 섹션의 세부 내용은 bullet list로 작성합니다.
- 불필요한 설명, 코드블록, 서문, 결론 멘트는 출력하지 않습니다.
- 출력은 반드시 Markdown만 사용합니다.
- 첫 줄은 "# 보고서 제목" 형식으로 작성합니다.

[작성일]
$today

[보고서 유형]
$meeting_focus

[사용자 요청]
$work_content

[한성대학교·사업단 참고자료]
$reference_section

[확인된 관련 규정 참고자료]
$rule_section

[선택된 Gmail 참고자료]
$email_section

[사용자 첨부자료]
$file_section

출력 구조는 반드시 [선택한 보고서 유형별 작성 지시]를 따르세요.
"""


DEFAULT_WEEKLY_REPORT_REQUEST_PROMPT = """아래는 한성대학교 SW중심대학 사업단 운영위원회 직원이 이번 주에 처리한 결재/업무 $approval_count건입니다.
이 내용을 바탕으로 주간보고서 초안을 작성해 주세요.

작성 방향:
- 이번 주 수행 업무를 예산·계획, 운영 규정·지침, 교육 정책·산학협력, 성과지표·환류 관점으로 묶어 정리
- 각 업무의 목적, 처리 내용, 현재 결재 상태, 후속 조치를 포함
- 운영위원회 보고에 적합한 공식 문체 사용
- 확인되지 않은 수치나 일정은 임의로 만들지 말고 확인 필요로 표시

[업무 분류 기준]
$category_text

[주간보고서 양식]
$template_text

[이번 주 처리 업무]
$weekly_work_text
"""


REPORT_TYPE_TEMPLATES = {
    "요약 보고서": {
        "purpose": "핵심 내용과 의사결정에 필요한 사항을 짧고 압축적으로 정리합니다.",
        "sections": [
            "## 1. 핵심 요약",
            "## 2. 주요 업무 및 결재 현황",
            "## 3. 주요 이슈",
            "## 4. 확인 필요 사항",
            "## 5. 후속 조치",
        ],
        "rules": [
            "각 항목은 2~4개 bullet로 간결하게 작성",
            "세부 설명보다 핵심 판단 정보 우선",
            "중복되는 업무는 하나의 묶음으로 요약",
        ],
    },
    "업무 정리 보고서": {
        "purpose": "담당자가 실제로 수행한 업무를 주간 업무보고 형식으로 정리합니다.",
        "sections": [
            "## 1. 금주 업무 총괄",
            "## 2. 업무별 처리 현황",
            "## 3. 결재 및 심의 진행 상황",
            "## 4. 산출물 및 첨부자료",
            "## 5. 차주 예정 업무",
        ],
        "rules": [
            "업무별 처리 결과와 후속 조치를 반드시 분리",
            "결재 상태를 결재대기/진행중/완료/반려로 명시",
            "업무분류 기준에 따라 묶어서 작성",
        ],
    },
    "회의 보고서": {
        "purpose": "운영위원회 또는 분과 회의에서 논의할 안건과 결정사항 중심으로 정리합니다.",
        "sections": [
            "## 1. 회의 개요",
            "## 2. 상정 안건",
            "## 3. 주요 논의 내용",
            "## 4. 결정 및 의결 사항",
            "## 5. 담당 부서별 후속 조치",
        ],
        "rules": [
            "회의 일시, 참석자 등 자료에 없는 정보는 확인 필요로 표시",
            "안건별 논의 내용과 결정 사항을 구분",
            "의결이 필요한 항목은 별도 표시",
        ],
    },
    "PPT 발표용 보고서": {
        "purpose": "PPT 슬라이드로 변환하기 좋은 발표 자료 구조로 정리합니다.",
        "sections": [
            "## 1. 발표 개요",
            "## 2. 핵심 메시지",
            "## 3. 주요 추진 내용",
            "## 4. 성과 및 이슈",
            "## 5. 향후 계획",
        ],
        "rules": [
            "각 bullet은 슬라이드에 들어갈 수 있게 짧게 작성",
            "한 섹션에 bullet 3~5개를 권장",
            "표현은 발표자가 읽기 쉬운 명사형 또는 간결한 문장으로 작성",
        ],
    },
    "검토 의견 보고서": {
        "purpose": "안건의 타당성, 쟁점, 리스크, 보완 의견을 검토자 관점으로 정리합니다.",
        "sections": [
            "## 1. 검토 배경",
            "## 2. 검토 대상",
            "## 3. 주요 쟁점",
            "## 4. 검토 의견",
            "## 5. 권고 및 보완 사항",
        ],
        "rules": [
            "찬반 판단보다 근거와 조건을 명확히 제시",
            "예산, 규정, KPI, 일정 리스크를 구분해 검토",
            "승인 전 확인이 필요한 사항을 반드시 포함",
        ],
    },
}


def report_type_instruction(meeting_focus: str) -> str:
    template = get_prompt_template(meeting_focus)
    sections = "\n".join(template["sections"])
    rules = "\n".join([f"- {rule}" for rule in template["rules"]])
    return f"""[선택한 보고서 유형별 작성 지시]
보고서 유형: {meeting_focus}
작성 목적: {template["purpose"]}

반드시 아래 섹션 제목과 순서를 사용하세요:
{sections}

유형별 작성 규칙:
{rules}
"""


def render_report_prompt(context: dict) -> str:
    template = Template(get_ai_prompt("base_report_prompt"))
    return template.safe_substitute(context)


def render_weekly_report_request(context: dict) -> str:
    template = Template(get_ai_prompt("weekly_report_request_prompt"))
    return template.safe_substitute(context)


def format_weekly_work_items(approvals: List[dict]) -> str:
    blocks = []
    for index, approval in enumerate(approvals, start=1):
        attachments = approval.get("attachments") or []
        attachment_names = ", ".join([file.get("name", "") for file in attachments if file.get("name")]) or "첨부 없음"
        blocks.append(
            f"""{index}. [{approval.get('category')}] {approval.get('title')}
- 처리일시: {approval.get('submittedAt')}
- 담당부서: {approval.get('department')}
- 업무분류: {approval.get('workType') or '확인 필요'}
- 산출물: {approval.get('outputType') or approval.get('category')}
- 결재상태: {approval.get('status')}
- 금액/예산: {approval.get('amount')}
- 주요내용: {approval.get('content')}
- 처리결과: {approval.get('result') or '확인 필요'}
- 후속조치: {approval.get('followUp') or '확인 필요'}
- 첨부서류: {attachment_names}"""
        )
    return "\n\n".join(blocks)


def format_work_categories() -> str:
    return "\n".join([f"- {item['name']}: {item['description']}" for item in WORK_CATEGORIES])


def format_weekly_template() -> str:
    return (
        f"{WEEKLY_REPORT_TEMPLATE['title']}\n"
        + "\n".join(WEEKLY_REPORT_TEMPLATE["sections"])
        + "\n\n작성 규칙:\n"
        + "\n".join([f"- {rule}" for rule in WEEKLY_REPORT_TEMPLATE["writingRules"]])
    )


seed_prompt_templates()
seed_ai_prompts()
ensure_ai_prompt_defaults()


PPT_NAVY = RGBColor(24, 50, 79)
PPT_BLUE = RGBColor(33, 90, 142)
PPT_TEAL = RGBColor(16, 127, 128)
PPT_BG = RGBColor(245, 247, 250)
PPT_TEXT = RGBColor(23, 32, 51)
PPT_MUTED = RGBColor(102, 112, 133)
PPT_WHITE = RGBColor(255, 255, 255)


def set_shape_fill(shape, color: RGBColor):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    shape.line.fill.background()


def add_template_text(slide, left, top, width, height, text, size=18, color=PPT_TEXT, bold=False, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0)
    frame.margin_right = Inches(0)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align is not None:
        paragraph.alignment = align
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def apply_slide_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = PPT_BG


def add_template_footer(slide, prs: Presentation, page_label: str):
    footer_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        prs.slide_height - Inches(0.36),
        prs.slide_width,
        Inches(0.36),
    )
    set_shape_fill(footer_bar, PPT_NAVY)
    add_template_text(
        slide,
        Inches(0.42),
        prs.slide_height - Inches(0.28),
        Inches(5.2),
        Inches(0.2),
        "Hansung SW-Centered University Committee",
        8,
        PPT_WHITE,
        False,
    )
    add_template_text(
        slide,
        prs.slide_width - Inches(1.15),
        prs.slide_height - Inches(0.28),
        Inches(0.72),
        Inches(0.2),
        page_label,
        8,
        PPT_WHITE,
        False,
        PP_ALIGN.RIGHT,
    )


def add_template_header(slide, prs: Presentation, title: str, page_index: int):
    apply_slide_background(slide)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.58))
    set_shape_fill(header, PPT_NAVY)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.58), prs.slide_width, Inches(0.06))
    set_shape_fill(accent, PPT_TEAL)
    add_template_text(slide, Inches(0.45), Inches(0.16), Inches(5.8), Inches(0.28), "SW중심대학사업단 운영위원회", 10, PPT_WHITE, True)
    add_template_text(slide, Inches(0.62), Inches(0.92), Inches(11.9), Inches(0.6), title[:70], 24, PPT_NAVY, True)
    add_template_footer(slide, prs, f"{page_index:02d}")


def add_bullet_block(slide, items: List[str], left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = None
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)

    display_items = items or ["내용 없음"]
    total_chars = sum(len(item) for item in display_items)
    if total_chars > 520 or len(display_items) >= 5:
        font_size = 13
        space_after = 3
        max_chars = 150
    elif total_chars > 360 or len(display_items) >= 4:
        font_size = 15
        space_after = 5
        max_chars = 180
    else:
        font_size = 17
        space_after = 8
        max_chars = 220

    for item_index, item in enumerate(display_items):
        paragraph = frame.paragraphs[0] if item_index == 0 else frame.add_paragraph()
        paragraph.text = item[:max_chars]
        paragraph.level = 0
        paragraph.space_after = Pt(space_after)
        paragraph.line_spacing = 1.05
        paragraph.font.name = "Malgun Gothic"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = PPT_TEXT


def add_title_template_slide(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide)
    left_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(1.0), prs.slide_height)
    set_shape_fill(left_band, PPT_NAVY)
    top_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), 0, prs.slide_width - Inches(1.0), Inches(0.18))
    set_shape_fill(top_accent, PPT_TEAL)
    add_template_text(slide, Inches(1.35), Inches(1.42), Inches(10.8), Inches(0.45), "SW중심대학사업단 운영위원회", 16, PPT_BLUE, True)
    add_template_text(slide, Inches(1.35), Inches(2.05), Inches(10.8), Inches(1.35), title[:80], 34, PPT_NAVY, True)
    add_template_text(slide, Inches(1.38), Inches(3.75), Inches(8.5), Inches(0.35), datetime.now().strftime("%Y년 %m월 %d일"), 15, PPT_MUTED)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.35), Inches(4.45), Inches(3.15), Inches(0.55))
    set_shape_fill(badge, PPT_BLUE)
    add_template_text(slide, Inches(1.58), Inches(4.61), Inches(2.7), Inches(0.2), "AI GENERATED DRAFT", 10, PPT_WHITE, True, PP_ALIGN.CENTER)
    add_template_footer(slide, prs, "00")


@app.post("/generate")
def generate_report(req: WorkRequest):
    today = datetime.now().strftime("%Y년 %m월 %d일")
    reference_section = search_reference_sources(req.work_content)
    rule_section = search_rules(req.work_content)
    email_section = format_email_sources(req.email_sources)
    file_section = format_uploaded_files(req.uploaded_files)
    file_parts = collect_gemini_file_parts(req.email_sources, req.uploaded_files)
    type_instruction = report_type_instruction(req.meeting_focus)

    prompt = render_report_prompt(
        {
            "type_instruction": type_instruction,
            "today": today,
            "meeting_focus": req.meeting_focus,
            "work_content": req.work_content or "직접 입력 없음",
            "reference_section": reference_section,
            "related": reference_section,
            "rule_section": rule_section,
            "email_section": email_section,
            "file_section": file_section,
        }
    )

    contents = [prompt, *file_parts] if file_parts else prompt
    response = gemini.models.generate_content(model="gemini-2.5-flash-lite", contents=contents)
    return {"draft": strip_code_fence(response.text)}


@app.post("/save-word")
def save_word(req: SaveRequest):
    media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    template_path = os.path.join(TEMPLATE_DIR, "weekly_report_template.docx")
    if os.path.exists(template_path):
        doc = Document(template_path)
        fill_weekly_report_template(doc, req.content)
    else:
        doc = Document()
        for line in markdown_lines(req.content):
            add_word_paragraph(doc, line)

    buffer = BytesIO()
    doc.save(buffer)
    file_bytes = buffer.getvalue()
    save_generated_file(req.filename, "word", media_type, file_bytes, req.content)
    return generated_file_response(file_bytes, req.filename, media_type)


@app.post("/save-ppt")
def save_ppt(req: SaveRequest):
    media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    title, sections = markdown_to_ppt_sections(req.content)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_template_slide(prs, title)
    slide_index = 1

    for section in sections:
        item_chunks = list(chunk_ppt_items(section["items"])) or [[]]
        for chunk_index, items in enumerate(item_chunks):
            slide_index += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            suffix = f" ({chunk_index + 1})" if len(item_chunks) > 1 else ""
            add_template_header(slide, prs, f"{section['title']}{suffix}", slide_index)
            content_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.62),
                Inches(1.75),
                Inches(12.1),
                Inches(4.95),
            )
            set_shape_fill(content_card, PPT_WHITE)
            add_bullet_block(slide, items, Inches(0.95), Inches(2.05), Inches(11.35), Inches(4.35))

    buffer = BytesIO()
    prs.save(buffer)
    file_bytes = buffer.getvalue()
    save_generated_file(req.filename, "ppt", media_type, file_bytes, req.content)
    return generated_file_response(file_bytes, req.filename, media_type)
